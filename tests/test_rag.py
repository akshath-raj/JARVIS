"""Tests for the document RAG: extraction, chunking, and the self-updating store.
Offline — a deterministic fake embedder stands in for Ollama/OpenAI."""
from __future__ import annotations

import hashlib

import numpy as np

from jarvis.rag import RAGPipeline, VectorStore
from jarvis.rag.chunker import chunk_sections
from jarvis.rag.extract import Section, extract_sections


class FakeEmbedder:
    key = "fake:test"

    def embed(self, texts):
        out = []
        for t in texts:
            h = hashlib.sha256(t.encode()).digest()
            v = np.frombuffer(h, dtype=np.uint8)[:8].astype(np.float32)
            n = np.linalg.norm(v) or 1.0
            out.append(v / n)
        return np.asarray(out, dtype=np.float32) if out else np.zeros((0, 8), np.float32)

    def embed_one(self, t):
        return self.embed([t])[0]


# ── extraction / chunking ────────────────────────────────────────────────────
def test_extract_markdown_sections(tmp_path):
    f = tmp_path / "n.md"
    f.write_text("# Alpha\nfirst body\n\n# Beta\nsecond body")
    secs = extract_sections(f)
    assert [s.title for s in secs] == ["Alpha", "Beta"]


def test_chunker_attaches_metadata(tmp_path):
    f = tmp_path / "n.md"
    secs = [Section(title="Intro", text="hello world. " * 200, page=None)]
    chunks = chunk_sections(f, secs, max_chars=300)
    assert len(chunks) > 1               # long section split into several chunks
    m = chunks[0].metadata
    assert m["filename"] == "n.md" and m["section"] == "Intro" and "cite" in m
    assert m["chunk_index"] == 0


def test_pptx_pages_and_docx_headings(tmp_path):
    from docx import Document
    from pptx import Presentation

    d = Document(); d.add_heading("H1", level=1); d.add_paragraph("body")
    d.save(str(tmp_path / "d.docx"))
    assert any(s.title == "H1" for s in extract_sections(tmp_path / "d.docx"))

    prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Slide A"; s.placeholders[1].text = "content"
    prs.save(str(tmp_path / "p.pptx"))
    secs = extract_sections(tmp_path / "p.pptx")
    assert secs[0].page == 1


# ── vector store ─────────────────────────────────────────────────────────────
def test_store_add_search_remove(tmp_path):
    emb = FakeEmbedder()
    store = VectorStore(str(tmp_path / "idx"), dim=8, embed_key=emb.key)

    class C:
        def __init__(self, text):
            self.text = text
            self.metadata = {"source": "a.txt", "cite": "a.txt"}

    chunks = [C("apple pie recipe"), C("rocket engine thrust")]
    store.add_document("a.txt", hash="h1", mtime=1.0, size=10,
                       chunks=chunks, vectors=emb.embed([c.text for c in chunks]))
    assert store.stats() == {"documents": 1, "chunks": 2}
    hits = store.search(emb.embed_one("apple pie recipe"), k=1)
    assert hits and hits[0]["text"] == "apple pie recipe"
    store.remove_document("a.txt")
    assert store.stats() == {"documents": 0, "chunks": 0}


def test_store_persists_and_reloads(tmp_path):
    emb = FakeEmbedder()
    d = str(tmp_path / "idx")

    class C:
        text = "persisted chunk"
        metadata = {"source": "a.txt", "cite": "a.txt"}

    VectorStore(d, dim=8, embed_key=emb.key).add_document(
        "a.txt", hash="h", mtime=1.0, size=1, chunks=[C()], vectors=emb.embed(["persisted chunk"]))
    reloaded = VectorStore(d, dim=8, embed_key=emb.key)
    assert reloaded.stats()["chunks"] == 1


def test_embed_key_change_invalidates_index(tmp_path):
    d = str(tmp_path / "idx")

    class C:
        text = "x"
        metadata = {"source": "a.txt", "cite": "a"}

    VectorStore(d, dim=8, embed_key="fake:v1").add_document(
        "a.txt", hash="h", mtime=1, size=1, chunks=[C()], vectors=FakeEmbedder().embed(["x"]))
    # a different embedding model must reset the index rather than mix spaces
    assert VectorStore(d, dim=8, embed_key="fake:v2").stats()["chunks"] == 0


# ── incremental pipeline sync (add / update / delete) ────────────────────────
def _pipeline(tmp_path, docs):
    store = VectorStore(str(tmp_path / "idx"), dim=8, embed_key="fake:test")
    return RAGPipeline(store, FakeEmbedder(), answer_model="x", dirs=[str(docs)]), store


def test_incremental_sync_add_update_delete(tmp_path):
    docs = tmp_path / "docs"; docs.mkdir()
    (docs / "a.txt").write_text("alpha content")
    (docs / "b.md").write_text("# H\nbeta content")
    rag, store = _pipeline(tmp_path, docs)

    s1 = rag.sync()
    assert s1["added"] == 2 and store.stats()["documents"] == 2

    # no changes → nothing re-embedded
    assert rag.sync() == {"added": 0, "updated": 0, "removed": 0, "unchanged": 2, "failed": 0}

    # edit one file → exactly one update
    import os, time
    (docs / "a.txt").write_text("alpha content changed substantially")
    os.utime(docs / "a.txt", (time.time() + 5, time.time() + 5))
    s3 = rag.sync()
    assert s3["updated"] == 1 and s3["unchanged"] == 1

    # delete a file → auto-removed from the index
    (docs / "b.md").unlink()
    s4 = rag.sync()
    assert s4["removed"] == 1 and store.stats()["documents"] == 1


def test_failed_extraction_is_cached_not_retried(tmp_path):
    docs = tmp_path / "docs"; docs.mkdir()
    (docs / "scanned.txt").write_text("")  # no extractable text → ExtractError
    rag, store = _pipeline(tmp_path, docs)
    s1 = rag.sync()
    assert s1["failed"] == 1
    # the failure is remembered, so the next scan skips it instead of re-trying
    s2 = rag.sync()
    assert s2["failed"] == 0 and s2["unchanged"] == 1


def test_scan_limit_spreads_work(tmp_path):
    docs = tmp_path / "docs"; docs.mkdir()
    for i in range(5):
        (docs / f"f{i}.txt").write_text(f"document number {i}")
    rag, store = _pipeline(tmp_path, docs)
    s = rag.sync(limit=2)          # only index 2 this cycle
    assert s["added"] == 2
    assert rag.sync(limit=2)["added"] == 2   # next cycle picks up 2 more
    assert rag.sync(limit=2)["added"] == 1   # and the last one


def test_ingest_single_file(tmp_path):
    docs = tmp_path / "docs"; docs.mkdir()
    rag, store = _pipeline(tmp_path, docs)
    f = docs / "new.txt"; f.write_text("freshly downloaded report")
    assert rag.ingest_file(str(f))["indexed"] is True
    assert store.stats()["documents"] == 1
    # second ingest is a no-op (already indexed, unchanged)
    assert rag.ingest_file(str(f))["indexed"] is False


# ── metadata enrichment (page / dates / location) ────────────────────────────
def test_chunk_metadata_has_dates_and_location(tmp_path):
    f = tmp_path / "report.txt"
    f.write_text("some real content here " * 20)
    chunks = chunk_sections(f, [Section(title="Intro", text="alpha " * 50, page=None)])
    m = chunks[0].metadata
    assert m["location"] == str(tmp_path)
    assert m["source"] == str(f)
    # created/modified are ISO dates (YYYY-MM-DD)
    assert len(m["created"]) == 10 and m["modified"].count("-") == 2


def test_pdf_page_number_flows_into_metadata(tmp_path):
    secs = [Section(title="Page 3", text="the quadratic formula is here", page=3)]
    chunks = chunk_sections(tmp_path / "math.pdf", secs)
    assert chunks[0].metadata["page"] == 3
    assert "p3" in chunks[0].metadata["cite"]


# ── description-based resolution, page targeting, summary ────────────────────
def _index_docs(tmp_path, files: dict):
    docs = tmp_path / "docs"; docs.mkdir()
    for name, text in files.items():
        (docs / name).write_text(text)
    rag, store = _pipeline(tmp_path, docs)
    rag.sync()
    return rag, store, docs


def test_resolve_document_by_filename_words(tmp_path):
    rag, _, docs = _index_docs(tmp_path, {
        "machine_learning_assignment.txt": "gradient descent and neural nets",
        "grocery_list.txt": "milk eggs bread",
    })
    doc = rag.resolve_document("my machine learning assignment")
    assert doc and doc["filename"] == "machine_learning_assignment.txt"
    assert doc["source"] == str(docs / "machine_learning_assignment.txt")


def test_resolve_document_prefers_exact_filename(tmp_path):
    # a caller that already knows the file (e.g. reviewing a doc it just opened) must
    # match it exactly, not get fuzzy-ranked onto a semantically noisier neighbour.
    rag, _, _ = _index_docs(tmp_path, {
        "Jake_s_Resume.txt": "brief resume",
        "Syllabus.txt": "course syllabus resume-building workshop pdf notes",
    })
    # bias the fuzzy ranker toward Syllabus, then confirm the exact name still wins
    rag.rank_documents = lambda desc, k=12: [
        {"source": "x", "filename": "Syllabus.txt", "score": 9.0},
        {"source": "y", "filename": "Jake_s_Resume.txt", "score": 0.1},
    ]
    doc = rag.resolve_document("Jake_s_Resume.txt")
    assert doc and doc["filename"] == "Jake_s_Resume.txt"           # exact name, not Syllabus
    assert rag.resolve_document("Jake_s_Resume")["filename"] == "Jake_s_Resume.txt"  # stem too


def test_related_documents_filename_match_beats_semantic_noise(tmp_path):
    # "resume" → the files whose NAME says resume, even if fuzzy content search buries
    # them under unrelated docs that happen to score higher.
    rag, _, _ = _index_docs(tmp_path, {
        "Jake_s_Resume.txt": "a",
        "my_resume.txt": "b",
        "lab_submission.txt": "c",
    })
    rag.rank_documents = lambda desc, k=40: [
        {"source": "l", "filename": "lab_submission.txt", "score": 9.0},   # semantic noise
        {"source": "j", "filename": "Jake_s_Resume.txt", "score": 1.5},
        {"source": "m", "filename": "my_resume.txt", "score": 1.4},
    ]
    names = {d["filename"] for d in rag.related_documents("resumes")}       # plural query
    assert names == {"Jake_s_Resume.txt", "my_resume.txt"}                  # noise excluded


def test_chunks_for_page_targeting(tmp_path):
    store = VectorStore(str(tmp_path / "idx"), dim=8, embed_key="fake:test")
    emb = FakeEmbedder()
    secs = [Section("Page 1", "first page text", page=1),
            Section("Page 2", "second page text", page=2)]
    chunks = chunk_sections(tmp_path / "deck.pdf", secs)
    store.add_document(str(tmp_path / "deck.pdf"), hash="h", mtime=1, size=1,
                       chunks=chunks, vectors=emb.embed([c.text for c in chunks]))
    only2 = store.chunks_for("deck.pdf", page=2)
    assert only2 and all(c["metadata"]["page"] == 2 for c in only2)
    assert store.documents()[0]["filename"] == "deck.pdf"


# ── disambiguation + on-screen subtopic explanation ──────────────────────────
def test_rank_documents_orders_and_scores(tmp_path):
    rag, _, _ = _index_docs(tmp_path, {
        "neural_networks.txt": "backpropagation gradient descent layers",
        "networks_notes.txt": "tcp ip routing subnets neural",
        "cooking.txt": "pasta tomato basil",
    })
    ranked = rag.rank_documents("neural network")
    assert ranked[0]["filename"] == "neural_networks.txt"
    assert all("score" in d for d in ranked)
    assert ranked[0]["score"] >= ranked[-1]["score"]


def test_related_documents_keeps_cluster_drops_noise(tmp_path):
    # the "open all docs on X" set: keep the top plus its close cluster, drop weak
    # keyword-only noise (below the absolute floor).
    rag, _, _ = _index_docs(tmp_path, {"a.txt": "anything"})
    ranked = [
        {"source": "top.pdf", "filename": "top.pdf", "score": 16.0},
        {"source": "b.pdf", "filename": "b.pdf", "score": 3.5},
        {"source": "c.pdf", "filename": "c.pdf", "score": 2.2},
        {"source": "noise.pdf", "filename": "noise.pdf", "score": 0.75},  # below floor
    ]
    rag.rank_documents = lambda desc, k=40: ranked
    rel = rag.related_documents("disaster")
    assert [d["filename"] for d in rel] == ["top.pdf", "b.pdf", "c.pdf"]


def test_related_documents_never_empty_when_there_is_a_hit(tmp_path):
    rag, _, _ = _index_docs(tmp_path, {"a.txt": "anything"})
    # even a single weak match (below the floor) is returned rather than nothing.
    rag.rank_documents = lambda desc, k=40: [{"source": "only.pdf", "filename": "only.pdf", "score": 0.4}]
    assert [d["filename"] for d in rag.related_documents("weak")] == ["only.pdf"]


def test_context_window_spans_pages_and_section(tmp_path):
    store = VectorStore(str(tmp_path / "idx"), dim=8, embed_key="fake:test")
    emb = FakeEmbedder()
    # a subtopic ("Entropy") running across pages 2-4, plus unrelated pages.
    secs = [
        Section("Page 1", "intro material", page=1),
        Section("Entropy", "entropy definition part one", page=2),
        Section("Entropy", "entropy worked example continues", page=3),
        Section("Entropy", "entropy summary and takeaways", page=4),
        Section("Page 5", "unrelated later material", page=5),
    ]
    chunks = chunk_sections(tmp_path / "thermo.pdf", secs)
    store.add_document(str(tmp_path / "thermo.pdf"), hash="h", mtime=1, size=1,
                       chunks=chunks, vectors=emb.embed([c.text for c in chunks]))
    # centre on the page-3 chunk; same-section chunks (2,3,4) must all be gathered.
    center = next(c["metadata"]["chunk_index"] for c in store.chunks_for("thermo.pdf", page=3))
    ctx = store.context_window("thermo.pdf", center, radius=0, section="Entropy")
    pages = sorted({c["metadata"]["page"] for c in ctx})
    assert pages == [2, 3, 4]        # whole subtopic, not just the visible page


def test_explain_topic_uses_located_subtopic(tmp_path, monkeypatch):
    rag, _, _ = _index_docs(tmp_path, {"os_notes.txt": "paging and virtual memory basics"})
    # stub the answering model so the test stays offline
    monkeypatch.setattr(rag, "_complete", lambda *a, **k: "Here is the explanation.")
    out = rag.explain_topic("explain this", document="os notes",
                            locator_text="virtual memory")
    assert out is not None and out["document"] == "os_notes.txt"
    assert out["answer"] == "Here is the explanation."


def test_explain_topic_falls_back_when_onscreen_doc_not_indexed(tmp_path, monkeypatch):
    rag, _, _ = _index_docs(tmp_path, {"operating_systems.txt": "paging virtual memory"})
    monkeypatch.setattr(rag, "_complete", lambda *a, **k: "ok")
    # the on-screen title matches the indexed file → explains from the document
    assert rag.explain_topic("x", document="operating systems", locator_text="paging") is not None
    # the live document's title matches NO indexed file → None, so the caller falls
    # back to plain screen vision instead of explaining the wrong file
    assert rag.explain_topic("x", document="quantum chemistry lab report", locator_text="q") is None


# ── summarising the document the user is looking at on screen ─────────────────
def test_summarize_located_whole_document(tmp_path, monkeypatch):
    rag, _, _ = _index_docs(tmp_path, {"os_notes.txt": "paging and virtual memory basics"})
    monkeypatch.setattr(rag, "_complete", lambda *a, **k: "The notes cover paging.")
    out = rag.summarize_located(document="os notes")
    assert out is not None and out["document"] == "os_notes.txt"
    assert "The notes cover paging." in out["answer"] and "os_notes.txt" in out["answer"]
    # whole-document scope → no page/subtopic narrowing
    assert out["page"] == 0 and out["subtopic"] == ""


def test_summarize_located_section_only_narrows_to_onscreen_subtopic(tmp_path, monkeypatch):
    store = VectorStore(str(tmp_path / "idx"), dim=8, embed_key="fake:test")
    emb = FakeEmbedder()
    secs = [
        Section("Page 1", "intro material", page=1),
        Section("Entropy", "entropy definition part one", page=2),
        Section("Entropy", "entropy worked example continues", page=3),
        Section("Entropy", "entropy summary and takeaways", page=4),
    ]
    chunks = chunk_sections(tmp_path / "thermo.pdf", secs)
    store.add_document(str(tmp_path / "thermo.pdf"), hash="h", mtime=1, size=1,
                       chunks=chunks, vectors=emb.embed([c.text for c in chunks]))
    rag = RAGPipeline(store, emb, answer_model="x")

    captured = {}

    def _fake_complete(prompt, **kwargs):
        captured["p"] = prompt
        return "Entropy summarised."

    monkeypatch.setattr(rag, "_complete", _fake_complete)
    out = rag.summarize_located(document="thermo", locator_text="entropy worked example",
                                page=3, section_only=True)
    assert out is not None
    assert out["subtopic"] == "Entropy" and out["page"] == 3
    # the located subtopic's content reached the model
    assert "entropy" in captured["p"].lower()


def test_summarize_located_falls_back_when_doc_not_indexed(tmp_path, monkeypatch):
    rag, _, _ = _index_docs(tmp_path, {"operating_systems.txt": "paging virtual memory"})
    monkeypatch.setattr(rag, "_complete", lambda *a, **k: "ok")
    # on-screen title matches the indexed file → summarises from the document
    assert rag.summarize_located(document="operating systems") is not None
    # unrelated on-screen title matches NO indexed file → None, so the caller falls
    # back to plain screen vision instead of summarising the wrong file
    assert rag.summarize_located(document="quantum chemistry lab report") is None
