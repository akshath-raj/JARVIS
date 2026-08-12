"""Extract a document into structured *sections* (title + text + page/slide), so
the chunker can keep semantically-related text together and tag each chunk with
where it came from.

  * PDF   → one section per page ("Page N")
  * DOCX  → sections split on Heading styles (fallback: whole doc)
  * PPTX  → one section per slide, titled by the slide title
  * TXT/MD→ sections split on markdown headings (fallback: paragraphs)

Dependency-light, reusing the same parsers as jarvis/documents.
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}


class ExtractError(RuntimeError):
    pass


@dataclass
class Section:
    title: str
    text: str
    page: int | None = None  # page (PDF) or slide (PPTX) number


def extract_sections(path: str | Path) -> list[Section]:
    p = Path(path)
    ext = p.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ExtractError(f"unsupported type: {ext or '(none)'}")
    try:
        if ext == ".pdf":
            secs = _pdf(p)
        elif ext == ".docx":
            secs = _docx(p)
        elif ext == ".pptx":
            secs = _pptx(p)
        else:
            secs = _text(p)
    except ExtractError:
        raise
    except Exception as e:  # noqa: BLE001
        raise ExtractError(f"couldn't read {p.name}: {e}") from e
    secs = [s for s in secs if s.text.strip()]
    if not secs:
        raise ExtractError(f"{p.name} has no extractable text (scanned image?)")
    return secs


def _pdf(p: Path) -> list[Section]:
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    out = []
    for i, page in enumerate(reader.pages, 1):
        out.append(Section(title=f"Page {i}", text=(page.extract_text() or "").strip(), page=i))
    return out


def _docx(p: Path) -> list[Section]:
    from docx import Document

    doc = Document(str(p))
    sections: list[Section] = []
    cur_title, cur_lines = "Introduction", []

    def flush():
        if cur_lines:
            sections.append(Section(title=cur_title, text="\n".join(cur_lines).strip()))

    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        txt = para.text.strip()
        if not txt:
            continue
        if style.startswith("heading") or style == "title":
            flush()
            cur_title, cur_lines = txt, []
        else:
            cur_lines.append(txt)
    flush()
    # include tables as their own section (often holds the substance)
    tbl_lines = []
    for table in doc.tables:
        for row in table.rows:
            tbl_lines.append(" | ".join(c.text for c in row.cells))
    if tbl_lines:
        sections.append(Section(title="Tables", text="\n".join(tbl_lines).strip()))
    return sections or [Section(title=p.stem, text="")]


def _pptx(p: Path) -> list[Section]:
    from pptx import Presentation

    prs = Presentation(str(p))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t:
                continue
            if not title and shape == getattr(slide.shapes, "title", None):
                title = t
            else:
                body.append(t)
        out.append(Section(title=title or f"Slide {i}", text="\n".join(body).strip(), page=i))
    return out


def _text(p: Path) -> list[Section]:
    raw = p.read_text(errors="replace")
    lines = raw.splitlines()
    is_md = p.suffix.lower() in (".md", ".markdown")
    sections: list[Section] = []
    cur_title, cur_lines = p.stem, []

    def flush():
        if cur_lines:
            sections.append(Section(title=cur_title, text="\n".join(cur_lines).strip()))

    heading = re.compile(r"^\s{0,3}(#{1,6})\s+(.*)$")
    for ln in lines:
        m = heading.match(ln) if is_md else None
        if m:
            flush()
            cur_title, cur_lines = m.group(2).strip(), []
        else:
            cur_lines.append(ln)
    flush()
    if not sections:  # plain text with no headings → whole file
        sections = [Section(title=p.stem, text=raw.strip())]
    return sections
