"""Read text out of the common document formats JARVIS downloads.

Kept dependency-light: pypdf (PDF), python-docx (DOCX), python-pptx (PPTX), plain
read for text/markdown. Used to feed the local model for the "explain it" step.
"""
from __future__ import annotations

from pathlib import Path


class DocumentError(RuntimeError):
    pass


def read_text(path: str | Path, *, max_chars: int = 60_000) -> str:
    """Extract plain text from a document. Truncated to `max_chars` for the LLM."""
    p = Path(path)
    if not p.exists():
        raise DocumentError(f"no such file: {p}")
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            text = _pdf(p)
        elif ext == ".docx":
            text = _docx(p)
        elif ext == ".pptx":
            text = _pptx(p)
        elif ext in (".txt", ".md", ".markdown", ".rtf"):
            text = p.read_text(errors="replace")
        else:
            raise DocumentError(f"unsupported document type: {ext or '(none)'}")
    except DocumentError:
        raise
    except Exception as e:  # noqa: BLE001 - surface any parser failure uniformly
        raise DocumentError(f"couldn't read {p.name}: {e}") from e
    text = text.strip()
    if not text:
        raise DocumentError(f"{p.name} has no extractable text (scanned image?)")
    return text[:max_chars]


def _pdf(p: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(p: Path) -> str:
    from docx import Document

    doc = Document(str(p))
    parts = [para.text for para in doc.paragraphs]
    for table in doc.tables:  # include table cell text (often holds the questions)
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)
