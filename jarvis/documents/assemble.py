"""Turn an AI's answer (markdown + fenced code) into a real answer file.

The AI returns the complete answer as text; we build the file the user actually
asked for: a Jupyter notebook, a Word doc, a PowerPoint, a Python script, or plain
markdown/text. Format is inferred from the user's words (or the assignment).
"""
from __future__ import annotations

import re
import time
from pathlib import Path

# user words / assignment hints → output extension
_FORMAT_HINTS: list[tuple[str, str]] = [
    (r"\b(ipynb|notebook|jupyter)\b", "ipynb"),
    (r"\b(power\s*point|pptx?|slides?|presentation|deck)\b", "pptx"),
    (r"\b(word|docx?|essay|report|document|write[- ]?up)\b", "docx"),
    (r"\b(python\s*(script|file)|\.py\b|code file)\b", "py"),
    (r"\b(markdown|\.md\b)\b", "md"),
    (r"\b(text file|\.txt\b|plain text)\b", "txt"),
]

_FENCE = re.compile(r"```([\w+-]*)\n(.*?)```", re.S)


def detect_format(user_text: str = "", assignment_text: str = "") -> str:
    """Pick the output extension from what the user said, else the assignment,
    else default to a Word document."""
    blob = f"{user_text}\n{assignment_text}".lower()
    for pattern, ext in _FORMAT_HINTS:
        if re.search(pattern, blob):
            return ext
    return "docx"


def _segments(answer: str) -> list[tuple[str, str, str]]:
    """Split answer into ordered (kind, lang, text) where kind is 'code'|'prose'."""
    segs: list[tuple[str, str, str]] = []
    pos = 0
    for m in _FENCE.finditer(answer):
        prose = answer[pos : m.start()].strip()
        if prose:
            segs.append(("prose", "", prose))
        segs.append(("code", (m.group(1) or "").strip(), m.group(2).rstrip("\n")))
        pos = m.end()
    tail = answer[pos:].strip()
    if tail:
        segs.append(("prose", "", tail))
    return segs or [("prose", "", answer.strip())]


def assemble(answer: str, out_format: str, out_dir: str, base_name: str = "answer") -> str:
    """Write `answer` as a file of `out_format` in `out_dir`; return the path."""
    out = Path(out_dir).expanduser()
    out.mkdir(parents=True, exist_ok=True)
    path = out / f"{base_name}-{time.strftime('%Y%m%d-%H%M%S')}.{out_format}"
    fmt = out_format.lower()
    if fmt == "ipynb":
        _ipynb(answer, path)
    elif fmt == "docx":
        _docx(answer, path)
    elif fmt == "pptx":
        _pptx(answer, path)
    elif fmt == "py":
        path.write_text(_code_only(answer) or answer)
    else:  # md / txt / anything else → raw
        path.write_text(answer)
    return str(path)


def _ipynb(answer: str, path: Path) -> None:
    import nbformat
    from nbformat.v4 import new_code_cell, new_markdown_cell, new_notebook

    nb = new_notebook()
    for kind, lang, text in _segments(answer):
        if kind == "code" and lang in ("", "python", "py", "python3"):
            nb.cells.append(new_code_cell(text))
        elif kind == "code":  # non-python fenced block → keep as markdown code
            nb.cells.append(new_markdown_cell(f"```{lang}\n{text}\n```"))
        else:
            nb.cells.append(new_markdown_cell(text))
    nbformat.write(nb, str(path))


def _docx(answer: str, path: Path) -> None:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    for kind, _lang, text in _segments(answer):
        if kind == "code":
            for line in text.splitlines() or [""]:
                run = doc.add_paragraph().add_run(line)
                run.font.name = "Menlo"
                run.font.size = Pt(9)
        else:
            for para in text.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para)
    doc.save(str(path))


def _pptx(answer: str, path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    # Split on markdown headings into slides; fallback = one slide.
    blocks = re.split(r"\n(?=#{1,3}\s)", answer.strip()) or [answer]
    layout = prs.slide_layouts[1]  # title + content
    for block in blocks:
        lines = block.strip().splitlines()
        if not lines:
            continue
        title = re.sub(r"^#{1,3}\s*", "", lines[0])[:120]
        body = "\n".join(lines[1:]).strip()
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if body and len(slide.placeholders) > 1:
            slide.placeholders[1].text = body[:2500]
    if not prs.slides:
        prs.slides.add_slide(layout).shapes.title.text = "Answer"
    prs.save(str(path))


def _code_only(answer: str) -> str:
    """Concatenate all python code fences (for a .py output)."""
    blocks = [t for k, lang, t in _segments(answer) if k == "code" and lang in ("", "python", "py", "python3")]
    return "\n\n".join(blocks)
